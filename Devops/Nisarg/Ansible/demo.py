from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Learn Ansible"
subtitle.text = "DevOps Training for College Seniors"

# Slide 1: Introduction to Ansible
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction to Ansible"
content.text = "• What is Ansible?\n• Advantages of Ansible"

# Slide 2: What is Ansible?
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "What is Ansible?"
content.text = "Ansible is an open-source IT automation tool.\nIt automates cloud provisioning, configuration management, application deployment, and many other IT needs."

# Slide 3: Advantages of Ansible
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Advantages of Ansible"
content.text = "• Simple to learn\n• Agentless architecture\n• Integrates with DevOps practices\n• Powerful and flexible"

# Slide 4: How Ansible Works
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "How Ansible Works"
content.text = "• Ansible Architecture\n• YAML Basics\n• Writing Ansible Playbooks\n• Installing Ansible on Windows using WSL\n• Ansible Project Structure"

# Slide 5: Ansible Architecture
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Ansible Architecture"
content.text = "Ansible uses a simple architecture:\n• Control Node\n• Managed Nodes\n• Modules\n• Plugins"

# Slide 6: YAML Basics
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "YAML Basics"
content.text = "YAML (YAML Ain't Markup Language) is a human-readable data format:\n• Used to write Ansible playbooks\n• Simple and easy to understand"

# Slide 7: Writing Ansible Playbooks
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Writing Ansible Playbooks"
content.text = "Playbooks are written in YAML:\n• Define tasks\n• Use modules\n• Manage configurations and deployments"

# Slide 8: Installing Ansible on Windows using WSL
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Installing Ansible on Windows using WSL"
content.text = "• Install Windows Subsystem for Linux (WSL)\n• Install a Linux distribution (e.g., Ubuntu)\n• Install Ansible using package manager"

# Slide 9: Ansible Project Structure
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Ansible Project Structure"
content.text = "• Inventory file\n• Configuration file\n• Playbooks\n• Roles\n• Modules"

# Slide 10: Jinja2 Templates and Variables
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Jinja2 Templates and Variables"
content.text = "• Introduction to Jinja2\n• Declaring Variables in Ansible\n• Using Templates in Playbooks\n• Rendering Variables in Configuration Files"

# Slide 11: Introduction to Jinja2
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction to Jinja2"
content.text = "Jinja2 is a templating engine for Python:\n• Used in Ansible for variable substitution\n• Allows dynamic generation of configuration files"

# Slide 12: Declaring Variables in Ansible
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Declaring Variables in Ansible"
content.text = "Variables can be declared in:\n• Inventory files\n• Playbooks\n• Roles\n• Extra vars (command line)"

# Slide 13: Using Templates in Playbooks
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Using Templates in Playbooks"
content.text = "Templates are used to create dynamic configuration files:\n• Use the `template` module\n• Reference Jinja2 templates\n• Pass variables to templates"

# Slide 14: Rendering Variables in Configuration Files
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Rendering Variables in Configuration Files"
content.text = "Variables are rendered in templates:\n• Use `{{ variable_name }}` syntax\n• Supports complex data structures\n• Apply Jinja2 filters for transformation"

# Slide 15: Nested Vars and Jinja2 Filters
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Nested Vars and Jinja2 Filters"
content.text = "• Understanding Nested Vars\n• Jinja2 Filters on Variables\n• Practical Examples"

# Slide 16: Understanding Nested Vars
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Understanding Nested Vars"
content.text = "Nested variables are structured data:\n• Dictionaries\n• Lists\n• Accessed using dot notation or bracket notation"

# Slide 17: Jinja2 Filters on Variables
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Jinja2 Filters on Variables"
content.text = "Filters modify variables:\n• Syntax: `{{ variable | filter }}`\n• Common filters: `default`, `upper`, `lower`, `replace`"

# Slide 18: Practical Examples
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Practical Examples"
content.text = "Example of nested vars and filters:\n• Define nested variables in playbook\n• Apply filters to manipulate data\n• Render filtered data in templates"

# Slide 19: Importing Vars File and Passing Variables at Runtime
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Importing Vars File and Passing Variables at Runtime"
content.text = "• Importing Vars Files\n• Passing Variables at Runtime\n• Hands-on Exercises"

# Slide 20: Importing Vars Files
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Importing Vars Files"
content.text = "Vars files allow separation of variables:\n• Define variables in separate YAML files\n• Use `vars_files` keyword to import\n• Example syntax and usage"

# Slide 21: Passing Variables at Runtime
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Passing Variables at Runtime"
content.text = "Variables can be passed at runtime:\n• Use `--extra-vars` option in CLI\n• Override default variables\n• Example command and usage"

# Slide 22: Hands-on Exercises
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Hands-on Exercises"
content.text = "Practice importing vars files and passing variables at runtime:\n• Create playbooks with vars files\n• Run playbooks with `--extra-vars`\n• Verify results and troubleshoot"

# Slide 23: Assessment and Feedback
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Assessment and Feedback"
content.text = "• Formative Assessment Activities\n• Summative Assessment Task\n• Reflection and Feedback"

# Slide 24: Formative Assessment Activities
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Formative Assessment Activities"
content.text = "• Short quizzes and exercises after each activity\n• Check understanding of key concepts"

# Slide 25: Summative Assessment Task
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
content.text = "Summative assessment task:\n• Create an Ansible playbook to automate a real-world task\n• Present the playbook and explain how it works"

# Slide 26: Reflection and Feedback
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Reflection and Feedback"
content.text = "• Reflection on learning outcomes\n• Feedback from students\n• Areas for improvement"

# Slide 27: Resources and Tools
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Resources and Tools"
content.text = "• Official Ansible Documentation\n• Ansible GitHub Repository\n• Online Tutorials and Courses\n• Books and Reference Materials"

# Slide 28: Closure
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Closure"
content.text = "• Recap of key points\n• Importance of automation in DevOps\n• Encouragement to explore Ansible further\n• Questions and answers"

# Save the presentation
prs.save('/mnt/data/Learn_Ansible_Presentation.pptx')
